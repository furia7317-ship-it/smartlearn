"""python-pptx 课件渲染服务。

支持**模板库**（用户在「生成 PPT 设置」里选择）：
- 内置主题（programmatic）：每套 = 背景色 + 标题/正文字体配色 + 顶部强调条，离线可用、无版权风险。
- 文件模板（.pptx 母版）：把任意 .pptx 丢进 `app/assets/ppt_templates/`，即作为母版被选用——
  渲染时清掉母版自带样张、保留其主题/母版，再把内容灌进去。键名 `file:<文件名>`。

`create_ppt(script, output_path, template)` 按 template 选择渲染方式；缺省用 academic。
"""

from __future__ import annotations

import asyncio
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

# ────────────────────────────── 内置主题 ──────────────────────────────
# 颜色 RRGGBB；font 同时作用于中英文（中文走 East-Asian typeface）。

THEMES: dict[str, dict[str, Any]] = {
    "academic": {
        "name": "学术简约", "desc": "白底深蓝标题，正式清爽，适合课堂讲义",
        "bg": "FFFFFF", "title": "1F3864", "body": "333333",
        "accent": "2E74B5", "subtitle": "7F7F7F", "font": "微软雅黑",
    },
    "tech": {
        "name": "科技蓝", "desc": "深蓝底亮青强调，科技感，适合前沿/工程主题",
        "bg": "0B1F3A", "title": "FFFFFF", "body": "D6E4F0",
        "accent": "21D4FD", "subtitle": "9FB3C8", "font": "微软雅黑",
    },
    "warm": {
        "name": "暖阳", "desc": "米色暖调，亲和柔和，适合入门/科普",
        "bg": "FFF8F0", "title": "B5651D", "body": "4A3B2A",
        "accent": "E8A04C", "subtitle": "A78B6F", "font": "微软雅黑",
    },
    "chalk": {
        "name": "板书", "desc": "墨绿黑板风，粉笔字，适合公式推导/演算",
        "bg": "1E3A2F", "title": "F5F5DC", "body": "E8F0E3",
        "accent": "F2C14E", "subtitle": "B7C9BC", "font": "楷体",
    },
    "business": {
        "name": "商务红", "desc": "浅灰底酒红强调，稳重正式，适合答辩/汇报",
        "bg": "F5F5F5", "title": "8B1E2D", "body": "2B2B2B",
        "accent": "C0392B", "subtitle": "8A8A8A", "font": "微软雅黑",
    },
    "fresh": {
        "name": "清新绿", "desc": "浅绿底墨绿标题，轻松明快，适合通识/兴趣课",
        "bg": "F2FAF5", "title": "1E5631", "body": "33453B",
        "accent": "3CB371", "subtitle": "789A86", "font": "微软雅黑",
    },
}

DEFAULT_TEMPLATE = "academic"

# 文件模板目录（用户可往里丢 .pptx 母版）。随 backend/app 一起分发。
TEMPLATE_DIR = Path(__file__).resolve().parent.parent.parent / "assets" / "ppt_templates"
_FILE_PREFIX = "file:"


def _file_templates() -> dict[str, Path]:
    """扫描 .pptx 文件模板：键 `file:<stem>` → 文件路径。"""
    out: dict[str, Path] = {}
    try:
        if TEMPLATE_DIR.is_dir():
            for p in sorted(TEMPLATE_DIR.glob("*.pptx")):
                if p.name.startswith("~$"):  # 跳过 Office 临时锁文件
                    continue
                out[f"{_FILE_PREFIX}{p.stem}"] = p
    except Exception:
        pass
    return out


def list_ppt_templates() -> list[dict[str, str]]:
    """供前端选择的模板列表（内置主题 + 文件母版）。"""
    items = [
        {"key": k, "name": v["name"], "desc": v["desc"], "kind": "builtin"}
        for k, v in THEMES.items()
    ]
    for key, path in _file_templates().items():
        items.append({"key": key, "name": path.stem, "desc": "自定义 .pptx 母版", "kind": "file"})
    return items


# ────────────────────────────── 渲染入口 ──────────────────────────────


async def create_ppt(
    script: dict[str, Any],
    output_path: str | Path,
    template: str | None = None,
) -> Path:
    """根据大纲脚本生成 PPT。template 覆盖 script["template"]；都没有则默认主题。"""
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    key = template or script.get("template") or DEFAULT_TEMPLATE

    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, _render, script, str(output_path), key)
    return output_path


def _render(script: dict[str, Any], output_path: str, key: str):
    files = _file_templates()
    if key in files:
        _render_from_file(script, output_path, files[key])
    else:
        theme = THEMES.get(key, THEMES[DEFAULT_TEMPLATE])
        _render_builtin(script, output_path, theme)


# ────────────────────────────── 样式工具 ──────────────────────────────


def _rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _set_ea_font(run, typeface: str) -> None:
    """给 run 设中文字体（python-pptx 未直接暴露，手写 XML），失败不影响渲染。"""
    try:
        rpr = run._r.get_or_add_rPr()
        ea = rpr.find(qn("a:ea"))
        if ea is None:
            ea = rpr.makeelement(qn("a:ea"), {})
            rpr.append(ea)
        ea.set("typeface", typeface)
    except Exception:
        pass


def _style_text_frame(
    tf, *, font: str | None, size: int, color: str | None, bold: bool = False
) -> None:
    rgb = _rgb(color) if color else None
    for para in tf.paragraphs:
        targets = para.runs or [para]
        for t in targets:
            try:
                f = t.font
                f.size = Pt(size)
                f.bold = bold
                if rgb is not None:
                    f.color.rgb = rgb
                if font:
                    f.name = font
            except Exception:
                pass
            if font and para.runs:
                _set_ea_font(t, font)


def _apply_background(slide, hex_str: str) -> None:
    try:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = _rgb(hex_str)
    except Exception:
        pass


def _add_accent_bar(slide, prs, hex_str: str) -> None:
    try:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), prs.slide_width, Pt(10))
        bar.fill.solid()
        bar.fill.fore_color.rgb = _rgb(hex_str)
        bar.line.fill.background()
        bar.shadow.inherit = False
    except Exception:
        pass


def _layout_for(prs, layout_name: str):
    try:
        if layout_name == "title":
            return prs.slide_layouts[0]
        if layout_name == "two_column" and len(prs.slide_layouts) > 3:
            return prs.slide_layouts[3]
        return prs.slide_layouts[1]
    except Exception:
        return prs.slide_layouts[0]


def _fill_content(slide, slide_data: dict[str, Any]):
    """把 content 填进正文占位符（含 two_column）。返回是否填充成功。"""
    content_items = slide_data.get("content", []) or []
    layout_name = slide_data.get("layout", "content")
    if layout_name == "two_column" and len(slide.placeholders) > 2:
        midpoint = max(1, (len(content_items) + 1) // 2)
        columns = [content_items[:midpoint], content_items[midpoint:]]
        for ph_index, items in zip((1, 2), columns, strict=False):
            try:
                tf = slide.placeholders[ph_index].text_frame
            except Exception:
                continue
            tf.clear()
            for i, item in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(item)
        return True
    if len(slide.placeholders) > 1:
        try:
            tf = slide.placeholders[1].text_frame
        except Exception:
            return False
        tf.clear()
        for i, item in enumerate(content_items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = str(item)
        return True
    return False


# ────────────────────────────── 内置主题渲染 ──────────────────────────────


def _render_builtin(script: dict[str, Any], output_path: str, theme: dict[str, Any]):
    prs = Presentation()
    font = theme["font"]

    # 标题页
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_background(slide, theme["bg"])
    _add_accent_bar(slide, prs, theme["accent"])
    if slide.shapes.title:
        slide.shapes.title.text = script.get("title", "课件")
        _style_text_frame(slide.shapes.title.text_frame, font=font, size=40, color=theme["title"], bold=True)
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "学枢 AI 生成"
        _style_text_frame(slide.placeholders[1].text_frame, font=font, size=18, color=theme["subtitle"])

    # 内容页
    for slide_data in script.get("slides", []):
        slide = prs.slides.add_slide(_layout_for(prs, slide_data.get("layout", "content")))
        _apply_background(slide, theme["bg"])
        _add_accent_bar(slide, prs, theme["accent"])
        if slide.shapes.title:
            slide.shapes.title.text = slide_data.get("title", "")
            _style_text_frame(slide.shapes.title.text_frame, font=font, size=28, color=theme["title"], bold=True)
        if _fill_content(slide, slide_data):
            # 给正文统一配色字体
            ph_idx = (1, 2) if slide_data.get("layout") == "two_column" else (1,)
            for i in ph_idx:
                if len(slide.placeholders) > i:
                    try:
                        _style_text_frame(slide.placeholders[i].text_frame, font=font, size=18, color=theme["body"])
                    except Exception:
                        pass

    prs.save(output_path)


# ────────────────────────────── 文件母版渲染 ──────────────────────────────


def _clear_slides(prs) -> None:
    """清空母版自带的样张幻灯片，保留 master/layout/theme。

    既要从 sldIdLst 摘掉引用，也要 drop 对应关系（否则样张 part 仍留在包里，新加幻灯片
    复用同名 slide1.xml → 保存时出现 Duplicate name 警告/重复 part）。
    """
    try:
        sld_id_lst = prs.slides._sldIdLst
        for sld in list(sld_id_lst):
            rid = sld.get(qn("r:id"))
            if rid:
                try:
                    prs.part.drop_rel(rid)
                except Exception:
                    pass
            sld_id_lst.remove(sld)
    except Exception:
        pass


def _render_from_file(script: dict[str, Any], output_path: str, base: Path):
    """用 .pptx 母版作主题：清掉样张，套母版主题灌内容；失败回落默认主题。"""
    try:
        prs = Presentation(str(base))
        _clear_slides(prs)

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title:
            slide.shapes.title.text = script.get("title", "课件")
        if len(slide.placeholders) > 1:
            try:
                slide.placeholders[1].text = "学枢 AI 生成"
            except Exception:
                pass

        for slide_data in script.get("slides", []):
            slide = prs.slides.add_slide(_layout_for(prs, slide_data.get("layout", "content")))
            if slide.shapes.title:
                slide.shapes.title.text = slide_data.get("title", "")
            _fill_content(slide, slide_data)

        prs.save(output_path)
    except Exception:
        # 母版损坏/不兼容 → 回落默认内置主题，保证一定出片
        _render_builtin(script, output_path, THEMES[DEFAULT_TEMPLATE])
