"""Safe, bounded extraction for files attached to tutor questions."""

from __future__ import annotations

import asyncio
import csv
import io
import re
from pathlib import Path
from typing import Any


MAX_ATTACHMENT_BYTES = 20 * 1024 * 1024
MAX_IMAGE_BYTES = 6 * 1024 * 1024
MAX_EXTRACTED_CHARS = 18_000

IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff"}
TEXT_EXTENSIONS = {".txt", ".md", ".markdown", ".csv", ".json", ".py", ".java", ".c", ".cpp", ".html", ".xml"}
DOCUMENT_EXTENSIONS = {".docx", ".pdf"}
PRESENTATION_EXTENSIONS = {".pptx"}
SPREADSHEET_EXTENSIONS = {".xlsx"}
SUPPORTED_EXTENSIONS = (
    IMAGE_EXTENSIONS
    | TEXT_EXTENSIONS
    | DOCUMENT_EXTENSIONS
    | PRESENTATION_EXTENSIONS
    | SPREADSHEET_EXTENSIONS
)


class AttachmentExtractionError(ValueError):
    """A public, actionable attachment validation/extraction error."""


def attachment_kind(filename: str, content_type: str = "") -> str:
    extension = Path(filename).suffix.lower()
    if extension in IMAGE_EXTENSIONS or content_type.lower().startswith("image/"):
        return "image"
    if extension == ".pdf":
        return "pdf"
    if extension == ".docx":
        return "document"
    if extension == ".pptx":
        return "presentation"
    if extension == ".xlsx":
        return "spreadsheet"
    if extension in TEXT_EXTENSIONS:
        return "text"
    raise AttachmentExtractionError(
        "不支持这种文件。可上传图片、PDF、Word(.docx)、PPT(.pptx)、Excel(.xlsx)、TXT、Markdown、CSV 或代码文件。"
    )


def _clean_text(value: str) -> str:
    text = (value or "").replace("\x00", "")
    text = re.sub(r"[ \t\u00a0]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    return text[:MAX_EXTRACTED_CHARS]


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "gb18030", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


def _extract_pdf(data: bytes) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:  # pragma: no cover - packaging guard
        raise AttachmentExtractionError("PDF 解析组件未安装，请安装 pypdf 后重试。") from exc
    try:
        reader = PdfReader(io.BytesIO(data))
        pages = []
        for index, page in enumerate(reader.pages[:80], 1):
            text = (page.extract_text() or "").strip()
            if text:
                pages.append(f"[第 {index} 页]\n{text}")
        return "\n\n".join(pages)
    except Exception as exc:  # noqa: BLE001
        raise AttachmentExtractionError(f"PDF 无法解析：{str(exc)[:160]}") from exc


def _extract_docx(data: bytes) -> str:
    try:
        from docx import Document

        document = Document(io.BytesIO(data))
        blocks = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
        for table in document.tables[:30]:
            for row in table.rows[:200]:
                values = [cell.text.strip() for cell in row.cells]
                if any(values):
                    blocks.append(" | ".join(values))
        return "\n\n".join(blocks)
    except Exception as exc:  # noqa: BLE001
        raise AttachmentExtractionError(f"Word 文档无法解析：{str(exc)[:160]}") from exc


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(data))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides[:120], 1):
            values = []
            for shape in slide.shapes:
                text = str(getattr(shape, "text", "") or "").strip()
                if text:
                    values.append(text)
            if values:
                slides.append(f"[第 {index} 页]\n" + "\n".join(values))
        return "\n\n".join(slides)
    except Exception as exc:  # noqa: BLE001
        raise AttachmentExtractionError(f"PPT 文档无法解析：{str(exc)[:160]}") from exc


def _extract_xlsx(data: bytes) -> str:
    try:
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        output: list[str] = []
        for sheet in workbook.worksheets[:20]:
            output.append(f"[工作表：{sheet.title}]")
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), 1):
                if row_index > 300:
                    break
                values = ["" if value is None else str(value) for value in row[:40]]
                if any(value.strip() for value in values):
                    output.append(" | ".join(values))
        return "\n".join(output)
    except Exception as exc:  # noqa: BLE001
        raise AttachmentExtractionError(f"Excel 文档无法解析：{str(exc)[:160]}") from exc


def extract_attachment(filename: str, content_type: str, data: bytes) -> dict[str, Any]:
    """Validate one upload and extract bounded text without persisting the file."""

    safe_name = Path(filename or "未命名文件").name[:180]
    kind = attachment_kind(safe_name, content_type)
    if not data:
        raise AttachmentExtractionError("文件为空，无法上传。")
    limit = MAX_IMAGE_BYTES if kind == "image" else MAX_ATTACHMENT_BYTES
    if len(data) > limit:
        label = "6MB" if kind == "image" else "20MB"
        raise AttachmentExtractionError(f"文件过大，{kind == 'image' and '图片' or '单个文件'}不能超过 {label}。")

    extension = Path(safe_name).suffix.lower()
    if extension and extension not in SUPPORTED_EXTENSIONS and kind != "image":
        attachment_kind(safe_name, content_type)  # raises the consistent message

    if kind == "image":
        text = ""
    elif kind == "pdf":
        text = _extract_pdf(data)
    elif kind == "document":
        text = _extract_docx(data)
    elif kind == "presentation":
        text = _extract_pptx(data)
    elif kind == "spreadsheet":
        text = _extract_xlsx(data)
    elif extension == ".csv":
        decoded = _decode_text(data)
        rows = list(csv.reader(io.StringIO(decoded)))[:500]
        text = "\n".join(" | ".join(row[:40]) for row in rows)
    else:
        text = _decode_text(data)

    cleaned = _clean_text(text)
    if kind != "image" and len(cleaned) < 2:
        raise AttachmentExtractionError("文件中没有提取到可读文字；如果是扫描版 PDF，请先转成图片再上传。")
    return {
        "name": safe_name,
        "kind": kind,
        "media_type": content_type or "application/octet-stream",
        "size": len(data),
        "extracted_text": cleaned,
    }


async def extract_tutor_attachment(filename: str, content_type: str, data: bytes) -> dict[str, Any]:
    """Recognize image/PDF uploads before they enter the tutor request.

    Images require the visual model because local text extraction cannot
    represent diagrams, formulas, or spatial relationships. PDFs use cloud OCR
    first and fall back to pypdf only when the cloud service is unavailable and
    the document already contains a usable text layer.
    """

    safe_name = Path(filename or "未命名文件").name[:180]
    kind = attachment_kind(safe_name, content_type)
    if not data:
        raise AttachmentExtractionError("文件为空，无法上传。")
    limit = MAX_IMAGE_BYTES if kind == "image" else MAX_ATTACHMENT_BYTES
    if len(data) > limit:
        label = "6MB" if kind == "image" else "20MB"
        raise AttachmentExtractionError(f"文件过大，{kind == 'image' and '图片' or '单个文件'}不能超过 {label}。")
    if kind == "image":
        payload = await asyncio.to_thread(extract_attachment, safe_name, content_type, data)
        try:
            from app.services.iflytek.recognition import recognize_image

            recognized = await recognize_image(data)
        except Exception as exc:  # noqa: BLE001
            from app.services.iflytek.recognition import AttachmentRecognitionError

            if isinstance(exc, AttachmentRecognitionError):
                raise AttachmentExtractionError(str(exc)) from exc
            raise AttachmentExtractionError(f"图片识别失败（{type(exc).__name__}），请稍后重试。") from exc
        payload.update(
            extracted_text=_clean_text(recognized),
            recognition_status="recognized",
            recognition_provider="iflytek-image-understanding",
            recognition_notice="图片已由讯飞图片理解完成识别，再交给智能教师解答。",
        )
        return payload

    if kind == "pdf":
        cloud_error = ""
        try:
            from app.services.iflytek.recognition import recognize_pdf

            recognized = await recognize_pdf(safe_name, data)
            return {
                "name": safe_name,
                "kind": kind,
                "media_type": content_type or "application/pdf",
                "size": len(data),
                "extracted_text": _clean_text(recognized),
                "recognition_status": "recognized",
                "recognition_provider": "iflytek-pdf-ocr",
                "recognition_notice": "PDF 已由讯飞 OCR 大模型完成识别，再交给智能教师解答。",
            }
        except Exception as exc:  # noqa: BLE001
            cloud_error = str(exc)[:180]
        try:
            payload = await asyncio.to_thread(extract_attachment, safe_name, content_type, data)
        except AttachmentExtractionError as exc:
            raise AttachmentExtractionError(
                f"PDF 云端识别未完成（{cloud_error or '服务暂不可用'}），且文档没有可用文字层。请稍后重试。"
            ) from exc
        payload.update(
            recognition_status="fallback",
            recognition_provider="local-pypdf",
            recognition_notice=f"PDF 云端识别未完成，已读取文档文字层：{cloud_error or '服务暂不可用'}",
        )
        return payload

    payload = await asyncio.to_thread(extract_attachment, safe_name, content_type, data)
    payload.update(
        recognition_status="parsed",
        recognition_provider="local-parser",
        recognition_notice="文件已完成本地结构化解析。",
    )
    return payload
