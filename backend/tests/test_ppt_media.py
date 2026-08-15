"""PPT media route regression tests."""

from __future__ import annotations

from typing import Any

import pytest


@pytest.mark.asyncio
async def test_create_ppt_wraps_slides_for_renderer(monkeypatch: pytest.MonkeyPatch) -> None:
    from app.routers import media
    from app.schemas.media import PptRequest

    captured: dict[str, Any] = {}

    def fake_create_task(student_id: str, kind: str, topic: str, script: dict[str, Any]) -> str:
        captured.update(
            {
                "student_id": student_id,
                "kind": kind,
                "topic": topic,
                "script": script,
            }
        )
        return "ppt-task-1"

    monkeypatch.setattr(media.media_task_manager, "create_task", fake_create_task)

    response = await media.create_ppt(
        PptRequest(
            topic="动态规划",
            student_id="test_student_001",
            slides=[
                {
                    "slide_num": 1,
                    "title": "核心思想",
                    "content": ["拆分重叠子问题", "保存中间结果"],
                    "layout": "content",
                }
            ],
        )
    )

    assert response.task_id == "ppt-task-1"
    assert captured == {
        "student_id": "test_student_001",
        "kind": "ppt",
        "topic": "动态规划",
        "script": {
            "title": "动态规划",
            "slides": [
                {
                    "slide_num": 1,
                    "title": "核心思想",
                    "content": ["拆分重叠子问题", "保存中间结果"],
                    "layout": "content",
                }
            ],
            "total_slides": 1,
        },
    }


@pytest.mark.asyncio
async def test_create_ppt_writes_two_column_slide_content(tmp_path) -> None:
    from pptx import Presentation

    from app.services.media.ppt import create_ppt

    output_path = tmp_path / "courseware.pptx"
    await create_ppt(
        {
            "title": "动态规划",
            "slides": [
                {
                    "slide_num": 1,
                    "title": "动态规划 vs 贪心",
                    "content": ["动态规划关注全局最优", "贪心关注局部最优"],
                    "layout": "two_column",
                }
            ],
        },
        output_path,
    )

    prs = Presentation(str(output_path))
    all_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

    assert "动态规划 vs 贪心" in all_text
    assert "动态规划关注全局最优" in all_text
    assert "贪心关注局部最优" in all_text
